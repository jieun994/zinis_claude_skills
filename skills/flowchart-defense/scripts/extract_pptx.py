#!/usr/bin/env python3
# PPTX 도형 텍스트 추출(보조). 흐름 연결(화살표)은 비전 판독이 정확하므로,
# 이 스크립트는 "텍스트 누락 방지"용으로만 쓴다. python-pptx가 있으면 도형별
# 좌표까지, 없으면 XML에서 텍스트만 뽑는다. 결과는 JSON으로 stdout에, 에러는 stderr로.
import sys
import os
import json
import zipfile
import re


def extract_with_pptx(path):
    """python-pptx로 슬라이드별 도형 텍스트 + 좌표 추출."""
    from pptx import Presentation  # type: ignore
    from pptx.util import Emu  # noqa: F401

    prs = Presentation(path)
    slides = []
    for si, slide in enumerate(prs.slides, 1):
        shapes = []
        for shape in slide.shapes:
            text = ""
            if shape.has_text_frame:
                text = "\n".join(p.text for p in shape.text_frame.paragraphs).strip()
            if not text:
                continue
            try:
                pos = {
                    "left": int(shape.left) if shape.left is not None else None,
                    "top": int(shape.top) if shape.top is not None else None,
                    "width": int(shape.width) if shape.width is not None else None,
                    "height": int(shape.height) if shape.height is not None else None,
                }
            except Exception:
                pos = {}
            shapes.append({"text": text, "pos": pos, "shape_type": str(shape.shape_type)})
        slides.append({"slide": si, "shapes": shapes})
    return slides


def extract_from_xml(path):
    """python-pptx가 없을 때: zip+XML에서 텍스트 런만 추출(좌표 없음)."""
    slides = []
    with zipfile.ZipFile(path) as z:
        names = sorted(
            (n for n in z.namelist() if re.match(r"ppt/slides/slide\d+\.xml$", n)),
            key=lambda n: int(re.search(r"(\d+)", n).group(1)),
        )
        for si, name in enumerate(names, 1):
            data = z.read(name).decode("utf-8", errors="replace")
            texts = re.findall(r"<a:t>(.*?)</a:t>", data, re.S)
            texts = [t.strip() for t in texts if t.strip()]
            slides.append({"slide": si, "shapes": [{"text": t} for t in texts]})
    return slides


def main():
    if len(sys.argv) < 2:
        print("usage: extract_pptx.py <file.pptx>", file=sys.stderr)
        sys.exit(2)
    path = sys.argv[1]
    if not os.path.isfile(path):
        print("file not found: %s" % path, file=sys.stderr)
        sys.exit(1)

    method = "python-pptx"
    try:
        slides = extract_with_pptx(path)
    except ImportError:
        method = "xml-fallback"
        try:
            slides = extract_from_xml(path)
        except Exception as e:  # noqa: BLE001
            print("xml extract failed: %s" % e, file=sys.stderr)
            sys.exit(1)
    except Exception as e:  # noqa: BLE001
        print("pptx extract failed: %s" % e, file=sys.stderr)
        sys.exit(1)

    out = {
        "method": method,
        "slide_count": len(slides),
        "slides": slides,
        "note": "화살표 연결관계는 이 텍스트만으로 알 수 없음 — 비전 판독으로 복원할 것.",
    }
    print(json.dumps(out, ensure_ascii=False, indent=2))


if __name__ == "__main__":
    main()
