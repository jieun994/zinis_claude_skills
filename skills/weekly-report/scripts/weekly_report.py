#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""weekly-report 스킬 핵심 스크립트.

두 가지 명령:
  read  <prev.pptx>            지난주 보고서에서 이어받을 값(누적 실적, 주차, 기간 등)을 JSON으로 출력
  build <data.json>           data.json 사양대로 보고서 PPT를 생성

진척률 6칸은 사용자가 손으로 하던 공식을 그대로 자동 계산한다(누적 % 4개 → 6칸):
  주간 계획 = 이번주 목표 - 지난주 실적
  주간 실적 = 이번주 실적 - 지난주 실적
  누적 계획 = 이번주 목표
  누적 실적 = 이번주 실적
  차주 주간 = 차주 목표 - 이번주 실적
  차주 누적 = 차주 목표

에러는 stderr로, 결과는 stdout(JSON)으로 출력한다.
"""
import sys
import io
import os
import re
import json
import copy

try:
    sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")
    sys.stderr = io.TextIOWrapper(sys.stderr.buffer, encoding="utf-8")
except Exception:
    pass

from pptx import Presentation


# ----------------------------------------------------------------------------
# 공통 헬퍼
# ----------------------------------------------------------------------------
def err(msg):
    print(msg, file=sys.stderr)


def iter_shapes(shapes):
    """그룹 안까지 모든 도형을 재귀적으로 순회."""
    for sh in shapes:
        yield sh
        if sh.shape_type == 6:  # GROUP
            try:
                for sub in iter_shapes(sh.shapes):
                    yield sub
            except Exception:
                pass


def all_shapes(prs):
    for idx, slide in enumerate(prs.slides):
        for sh in iter_shapes(slide.shapes):
            yield idx, slide, sh


def table_of(sh):
    return sh.table if getattr(sh, "has_table", False) else None


def cell_text(cell):
    return cell.text.replace("\x0b", "\n")


def set_cell_text(cell, text):
    """셀 텍스트를 교체하되 첫 run의 서식(폰트/크기/색)을 보존한다.
    여러 줄은 '\n'으로 구분하며 같은 서식으로 단락을 만든다."""
    if text is None:
        return
    from pptx.oxml.ns import qn
    text = str(text)
    tf = cell.text_frame
    paras = tf.paragraphs
    # 첫 단락의 첫 run을 서식 템플릿으로 사용
    template_run = None
    if paras and paras[0].runs:
        template_run = paras[0].runs[0]

    lines = text.split("\n")

    # 첫 단락의 run과 줄바꿈(<a:br/>)을 모두 제거 — 잔여 br이 빈 줄로 남는 것 방지
    p0 = paras[0]
    for child in list(p0._p):
        if child.tag in (qn("a:r"), qn("a:br")):
            p0._p.remove(child)
    run0 = p0.add_run()
    run0.text = lines[0]
    if template_run is not None:
        _copy_run_format(template_run, run0)

    # 나머지 단락 제거
    for extra in list(tf.paragraphs[1:]):
        extra._p.getparent().remove(extra._p)

    # 추가 줄
    for line in lines[1:]:
        p = tf.add_paragraph()
        r = p.add_run()
        r.text = line
        if template_run is not None:
            _copy_run_format(template_run, r)


def _copy_run_format(src, dst):
    """run 서식(bold/size/color/이름) 복사 — XML rPr 통째로 복제."""
    try:
        from pptx.oxml.ns import qn
        src_rpr = src._r.find(qn("a:rPr"))
        if src_rpr is not None:
            new_rpr = copy.deepcopy(src_rpr)
            old = dst._r.find(qn("a:rPr"))
            if old is not None:
                dst._r.remove(old)
            dst._r.insert(0, new_rpr)
    except Exception:
        pass


def set_textbox_text(sh, text):
    """텍스트박스 전체를 한 줄 텍스트로 교체(첫 run 서식 보존)."""
    if not sh.has_text_frame:
        return
    from pptx.oxml.ns import qn
    tf = sh.text_frame
    p0 = tf.paragraphs[0]
    template_run = p0.runs[0] if p0.runs else None
    for child in list(p0._p):
        if child.tag in (qn("a:r"), qn("a:br")):
            p0._p.remove(child)
    run0 = p0.add_run()
    run0.text = text
    if template_run is not None:
        _copy_run_format(template_run, run0)
    for extra in list(tf.paragraphs[1:]):
        extra._p.getparent().remove(extra._p)


def fmt_pct(v):
    """숫자를 보고서 표기(예: 5.2%)로. None/'-'는 '-'."""
    if v is None or v == "-" or v == "":
        return "-"
    try:
        f = float(v)
    except (TypeError, ValueError):
        return str(v)
    # 원본 양식과 동일하게 항상 소수 첫째자리 (예: 5.2%, 7.0%, 11.8%)
    return ("%.1f" % f) + "%"


# ----------------------------------------------------------------------------
# 진척률 계산 (사용자 공식)
# ----------------------------------------------------------------------------
def compute_progress(prev_actual, this_plan, this_actual, next_plan):
    """누적 % 4개 → 표 6칸. 입력 중 하나라도 None이면 해당 계산은 '-'."""
    def sub(a, b):
        if a is None or b is None:
            return None
        return round(float(a) - float(b), 1)

    return {
        "주간_계획": sub(this_plan, prev_actual),
        "주간_실적": sub(this_actual, prev_actual),
        "누적_계획": this_plan,
        "누적_실적": this_actual,
        "차주_주간": sub(next_plan, this_actual),
        "차주_누적": next_plan,
    }


# ----------------------------------------------------------------------------
# 표 식별 (내용 시그니처 기반 — 슬라이드 순서가 바뀌어도 안전)
# ----------------------------------------------------------------------------
def find_tables(prs, predicate):
    out = []
    for idx, slide, sh in all_shapes(prs):
        tbl = table_of(sh)
        if tbl is not None and predicate(tbl):
            out.append((idx, slide, sh, tbl))
    return out


def is_progress_table(tbl):
    """전체/업무별 진척률표: 헤더 R0에 '금주 진척률'과 '차주 계획 진척률' 포함."""
    try:
        row0 = [tbl.cell(0, c).text for c in range(len(tbl.columns))]
    except Exception:
        return False
    joined = " ".join(row0)
    return ("금주 진척률" in joined) and ("차주 계획" in joined)


def is_weekly_box(tbl):
    """슬라이드2 진척률 박스 2x2: R0 = ['금주 실적', '금주 누적'|'전체 누적']."""
    try:
        if len(tbl.rows) != 2 or len(tbl.columns) != 2:
            return False
        return "금주 실적" in tbl.cell(0, 0).text
    except Exception:
        return False


def is_task_table(tbl):
    """주간 실적/차주 계획표: R0[0]=='업무구분' + 5열 + '관련 산출물' 헤더.
    (월간 성과/목표표도 '업무구분'으로 시작하므로 열 구성으로 구분)"""
    try:
        if tbl.cell(0, 0).text.strip() != "업무구분":
            return False
        if len(tbl.columns) < 5:
            return False
        row0 = " ".join(tbl.cell(0, c).text for c in range(len(tbl.columns)))
        return "산출물" in row0
    except Exception:
        return False


def is_issue_table(tbl):
    """주요 이슈표: R0 == ['구분','내용'] 이고 '요청 사항' 행 존재."""
    try:
        if tbl.cell(0, 0).text.strip() != "구분":
            return False
        col0 = [tbl.cell(r, 0).text.strip() for r in range(len(tbl.rows))]
        return "요청 사항" in col0
    except Exception:
        return False


def is_delay_table(tbl):
    try:
        return tbl.cell(0, 0).text.strip() == "지연 사유"
    except Exception:
        return False


# ----------------------------------------------------------------------------
# READ: 지난주 보고서에서 이어받을 값 추출
# ----------------------------------------------------------------------------
def cmd_read(path):
    if not os.path.isfile(path):
        err("파일을 찾을 수 없습니다: %s" % path)
        return 2
    prs = Presentation(path)
    out = {"source": os.path.basename(path), "mode": None, "week_label": None,
           "report_date": None, "period": None, "cumulative_actual": {}, "task_units": []}

    # 표지: W1/W2, 주차, 날짜
    for idx, slide, sh in all_shapes(prs):
        if sh.has_text_frame:
            t = sh.text_frame.text
            m = re.search(r"보고서\s*\((W\d)\)\s*[–-]\s*([^\n]+)", t)
            if m:
                out["mode"] = m.group(1)
                out["week_label"] = m.group(2).strip()
            if re.fullmatch(r"20\d{2}\.\s?\d{1,2}\.\s?\d{1,2}\.?", t.strip()):
                out["report_date"] = t.strip()

    # 전체/업무별 진척률표에서 '누적 실적'(=다음주 '지난주 실적')을 구분별로 추출
    # 컬럼 의미: R2 = ['','계획','실적','계획','실적','주간','누적'] → 누적실적 = col 4
    for idx, slide, sh, tbl in find_tables(prs, is_progress_table):
        ncols = len(tbl.columns)
        for r in range(3, len(tbl.rows)):
            gubun = tbl.cell(r, 0).text.strip()
            if not gubun:
                continue
            val = tbl.cell(r, 4).text.strip() if ncols > 4 else ""
            num = _parse_pct(val)
            if gubun and gubun not in out["cumulative_actual"]:
                out["cumulative_actual"][gubun] = num

    # 기간 텍스트
    for idx, slide, sh in all_shapes(prs):
        if sh.has_text_frame and sh.text_frame.text.strip().startswith("기간 :"):
            out["period"] = sh.text_frame.text.strip()
            break

    # 업무구분 목록(주간 실적표 기준)
    seen = []
    for idx, slide, sh, tbl in find_tables(prs, is_task_table):
        for r in range(1, len(tbl.rows)):
            g = tbl.cell(r, 0).text.strip()
            if g and g not in seen:
                seen.append(g)
    out["task_units"] = seen

    print(json.dumps(out, ensure_ascii=False, indent=2))
    return 0


def _parse_pct(s):
    if s is None:
        return None
    s = s.strip().replace("%", "")
    if s in ("", "-"):
        return None
    try:
        return float(s)
    except ValueError:
        return None


# ----------------------------------------------------------------------------
# BUILD: data.json 사양대로 보고서 생성
# ----------------------------------------------------------------------------
def fill_progress_tables(prs, progress_inputs):
    """progress_inputs: {구분: {prev_actual, this_plan, this_actual, next_plan}}.
    전체 진척률표 + 업무별 진척률표 모두 채운다.
    표 6칸 컬럼 인덱스: 1=주간계획 2=주간실적 3=누적계획 4=누적실적 5=차주주간 6=차주누적"""
    for idx, slide, sh, tbl in find_tables(prs, is_progress_table):
        for r in range(3, len(tbl.rows)):
            gubun = tbl.cell(r, 0).text.strip()
            data = _match_gubun(progress_inputs, gubun)
            if not data:
                continue
            cells = compute_progress(data.get("prev_actual"), data.get("this_plan"),
                                     data.get("this_actual"), data.get("next_plan"))
            order = ["주간_계획", "주간_실적", "누적_계획", "누적_실적", "차주_주간", "차주_누적"]
            for ci, key in enumerate(order, start=1):
                if ci < len(tbl.columns):
                    set_cell_text(tbl.cell(r, ci), fmt_pct(cells[key]))


def _match_gubun(inputs, gubun):
    """표의 구분명과 입력 키를 느슨하게 매칭(괄호/공백 차이 허용)."""
    if not gubun:
        return None
    norm = lambda s: re.sub(r"\s|\(.*?\)", "", s)
    g = norm(gubun)
    for k, v in inputs.items():
        if norm(k) == g or norm(k) in g or g in norm(k):
            return v
    return None


def fill_weekly_box(prs, total_inputs):
    """슬라이드2 진척률 박스: [금주 실적, 누적]을 전체 주간실적/누적실적으로."""
    if not total_inputs:
        return
    cells = compute_progress(total_inputs.get("prev_actual"), total_inputs.get("this_plan"),
                             total_inputs.get("this_actual"), total_inputs.get("next_plan"))
    for idx, slide, sh, tbl in find_tables(prs, is_weekly_box):
        set_cell_text(tbl.cell(1, 0), fmt_pct(cells["주간_실적"]))
        set_cell_text(tbl.cell(1, 1), fmt_pct(cells["누적_실적"]))


def fill_task_tables(prs, tasks, week_no, this_range, next_week_no, next_range):
    """주간 실적/차주 계획표 채우기.
    컬럼: 0=업무구분 1=금주 실적 2=산출물 3=차주 계획 4=산출물"""
    for idx, slide, sh, tbl in find_tables(prs, is_task_table):
        # 헤더 라벨 갱신
        if week_no:
            set_cell_text(tbl.cell(0, 1), "W%s 업무 진행\n(%s)" % (week_no, this_range or ""))
        if next_week_no:
            set_cell_text(tbl.cell(0, 3), "W%s 업무 계획\n(%s)" % (next_week_no, next_range or ""))
        for r in range(1, len(tbl.rows)):
            gubun = tbl.cell(r, 0).text.strip()
            data = tasks.get(gubun) or _match_task(tasks, gubun)
            if not data:
                continue
            if "this_done" in data:
                set_cell_text(tbl.cell(r, 1), data.get("this_done", ""))
            if "this_output" in data and len(tbl.columns) > 2:
                set_cell_text(tbl.cell(r, 2), data.get("this_output", ""))
            if "next_plan" in data and len(tbl.columns) > 3:
                set_cell_text(tbl.cell(r, 3), data.get("next_plan", ""))
            if "next_output" in data and len(tbl.columns) > 4:
                set_cell_text(tbl.cell(r, 4), data.get("next_output", ""))


def _match_task(tasks, gubun):
    norm = lambda s: re.sub(r"\s|\(.*?\)", "", s)
    g = norm(gubun)
    for k, v in tasks.items():
        if norm(k) == g:
            return v
    return None


def fill_issue_table(prs, issues):
    if not issues:
        return
    for idx, slide, sh, tbl in find_tables(prs, is_issue_table):
        for r in range(1, len(tbl.rows)):
            label = tbl.cell(r, 0).text.strip()
            if label in issues:
                set_cell_text(tbl.cell(r, 1), issues[label])


def fill_delay_table(prs, delay):
    if not delay:
        return
    for idx, slide, sh, tbl in find_tables(prs, is_delay_table):
        if "지연 사유" in delay:
            set_cell_text(tbl.cell(0, 1), delay["지연 사유"])
        if "캐치업 방안" in delay and len(tbl.rows) > 1:
            set_cell_text(tbl.cell(1, 1), delay["캐치업 방안"])


def _set_para_text(para, newtext):
    """단락의 첫 run 텍스트만 교체하고 서식(rPr) 보존, 나머지 run·줄바꿈 제거."""
    from pptx.oxml.ns import qn
    template_run = para.runs[0] if para.runs else None
    for child in list(para._p):
        if child.tag in (qn("a:r"), qn("a:br")):
            para._p.remove(child)
    run = para.add_run()
    run.text = newtext
    if template_run is not None:
        _copy_run_format(template_run, run)


def update_cover_and_dates(prs, mode, week_label, report_date, period, actual_basis):
    for idx, slide, sh in all_shapes(prs):
        if not sh.has_text_frame:
            continue
        tf = sh.text_frame
        t = tf.text
        # 표지 제목: '보고서 (Wn)' 줄만 단락 단위로 교체(다른 줄·서식 보존)
        if "보고서 (" in t:
            for para in tf.paragraphs:
                pt = "".join(r.text for r in para.runs)
                if "보고서 (" not in pt:
                    continue
                newp = re.sub(r"\(W\d\)", "(%s)" % mode, pt)
                if week_label:
                    newp = re.sub(r"([–-]\s*).+$", r"\g<1>%s" % week_label, newp)
                _set_para_text(para, newp)
            continue
        # 보고 날짜 (실제 날짜 또는 빈 템플릿의 20〇〇.〇〇.〇〇)
        if report_date and re.fullmatch(r"20[\d〇]{2}\.\s?[\d〇]{1,2}\.\s?[\d〇]{1,2}\.?", t.strip()):
            set_textbox_text(sh, report_date)
            continue
        # 기간
        if period and t.strip().startswith("기간"):
            set_textbox_text(sh, period)
            continue
        # 실적 기준
        if actual_basis and t.strip().startswith("실적 기준"):
            set_textbox_text(sh, "실적 기준 : %s" % actual_basis)
            continue


# ----------------------------------------------------------------------------
# W2 월간 결산 슬라이드 삽입 + 채우기
# ----------------------------------------------------------------------------
def has_monthly_section(prs):
    for idx, slide, sh in all_shapes(prs):
        if sh.has_text_frame and "진척 상황 보고" in sh.text_frame.text:
            return True
    return False


def base_project_title(prs):
    from collections import Counter
    c = Counter()
    for idx, slide, sh in all_shapes(prs):
        if sh.has_text_frame:
            t = sh.text_frame.text.strip()
            if "프로젝트" in t and 0 < len(t) < 60:
                c[t] += 1
    return c.most_common(1)[0][0] if c else None


def _pick_layout(prs):
    for layout in prs.slide_layouts:
        if "제목 및 내용" in layout.name:
            return layout
    return prs.slide_layouts[min(1, len(prs.slide_layouts) - 1)]


def append_monthly_slides(prs, asset_path, project_title):
    """asset_path의 4개 슬라이드를 prs에 복사하고 #Appendix 앞으로 이동."""
    if not os.path.isfile(asset_path):
        err("월간 결산 템플릿(assets/w2_monthly.pptx)을 찾을 수 없습니다.")
        return []
    src = Presentation(asset_path)
    layout = _pick_layout(prs)
    sldIdLst = prs.slides._sldIdLst
    before_ids = list(sldIdLst)

    new_slides = []
    for src_slide in src.slides:
        dst = prs.slides.add_slide(layout)
        for shp in list(dst.shapes):  # 레이아웃 기본 placeholder 제거
            shp._element.getparent().remove(shp._element)
        for shp in src_slide.shapes:
            dst.shapes._spTree.append(copy.deepcopy(shp._element))
        new_slides.append(dst)
        # 제목(프로젝트명) 채우기
        if project_title:
            for sh in dst.shapes:
                if sh.has_text_frame and sh.text_frame.text.strip() == "":
                    # 제목 placeholder 추정: 비어있고 type이 title 계열
                    try:
                        if sh.is_placeholder and sh.placeholder_format.idx == 0:
                            set_textbox_text(sh, project_title)
                    except Exception:
                        pass

    # 새로 추가된 sldId들을 #Appendix 앞으로 이동
    after_ids = list(sldIdLst)
    appended = [x for x in after_ids if x not in before_ids]
    # appendix 위치 찾기
    appendix_idx = None
    for i, slide in enumerate(prs.slides):
        txt = " ".join(sh.text_frame.text for sh in slide.shapes if sh.has_text_frame)
        if "Appendix" in txt:
            appendix_idx = i
            break
    if appendix_idx is not None:
        for el in appended:
            sldIdLst.remove(el)
        anchor = list(sldIdLst)[appendix_idx]
        for el in appended:
            anchor.addprevious(el)
    return new_slides


def is_month_progress(tbl):
    try:
        row0 = " ".join(tbl.cell(0, c).text for c in range(len(tbl.columns)))
        return "당월 진척률" in row0
    except Exception:
        return False


def is_perf_table(tbl):
    try:
        return "주요 성과" in tbl.cell(0, 1).text
    except Exception:
        return False


def is_goal_table(tbl):
    try:
        return "주요 목표" in tbl.cell(0, 1).text
    except Exception:
        return False


def is_risk_table(tbl):
    try:
        return tbl.cell(0, 0).text.strip() == "이슈 및 리스크"
    except Exception:
        return False


def fill_monthly(prs, monthly):
    """monthly = {month_label, next_month, period, plan_basis,
                  progress:{구분:[6값]}, performance:{구분:{text,note}},
                  issues:[{issue,impact,action}], goals:{구분:{text,note}}}"""
    if not monthly:
        return
    ml = monthly.get("month_label")        # 예: '6월'
    nm = monthly.get("next_month")          # 예: '7월'
    # 섹션 제목/헤더의 〇월 치환
    for idx, slide, sh in all_shapes(prs):
        if sh.has_text_frame:
            t = sh.text_frame.text
            if "〇월" in t:
                new = t.replace("〇월 진척 상황", (ml or "") + " 진척 상황")
                set_textbox_text(sh, new)
            elif t.strip() == "기간 :" and monthly.get("period"):
                set_textbox_text(sh, "기간 : " + monthly["period"])

    prog = monthly.get("progress", {})
    for idx, slide, sh, tbl in find_tables(prs, is_month_progress):
        # 헤더 〇월 치환
        for c in range(len(tbl.columns)):
            if "〇월" in tbl.cell(0, c).text and nm:
                set_cell_text(tbl.cell(0, c), tbl.cell(0, c).text.replace("〇월", nm))
        for r in range(3, len(tbl.rows)):
            g = tbl.cell(r, 0).text.strip()
            vals = prog.get(g) or _match_gubun(prog, g)
            if not vals:
                continue
            for ci, v in enumerate(vals[:6], start=1):
                if ci < len(tbl.columns):
                    set_cell_text(tbl.cell(r, ci), fmt_pct(v) if v not in (None, "-", "") else "-")

    perf = monthly.get("performance", {})
    for idx, slide, sh, tbl in find_tables(prs, is_perf_table):
        if ml:
            set_cell_text(tbl.cell(0, 1), ml + " 주요 성과")
        for r in range(1, len(tbl.rows)):
            g = tbl.cell(r, 0).text.strip()
            d = perf.get(g) or _match_task(perf, g)
            if not d:
                continue
            set_cell_text(tbl.cell(r, 1), d.get("text", ""))
            if len(tbl.columns) > 2:
                set_cell_text(tbl.cell(r, 2), d.get("note", ""))

    goals = monthly.get("goals", {})
    for idx, slide, sh, tbl in find_tables(prs, is_goal_table):
        if nm:
            set_cell_text(tbl.cell(0, 1), nm + " 주요 목표")
        for r in range(1, len(tbl.rows)):
            g = tbl.cell(r, 0).text.strip()
            d = goals.get(g) or _match_task(goals, g)
            if not d:
                continue
            set_cell_text(tbl.cell(r, 1), d.get("text", ""))
            if len(tbl.columns) > 2:
                set_cell_text(tbl.cell(r, 2), d.get("note", ""))

    issues = monthly.get("issues", [])
    for idx, slide, sh, tbl in find_tables(prs, is_risk_table):
        for r in range(1, len(tbl.rows)):
            if r - 1 < len(issues):
                it = issues[r - 1]
                set_cell_text(tbl.cell(r, 0), it.get("issue", ""))
                set_cell_text(tbl.cell(r, 1), it.get("impact", ""))
                if len(tbl.columns) > 2:
                    set_cell_text(tbl.cell(r, 2), it.get("action", ""))


def cmd_build(data_path):
    if not os.path.isfile(data_path):
        err("data.json을 찾을 수 없습니다: %s" % data_path)
        return 2
    with open(data_path, "r", encoding="utf-8") as f:
        data = json.load(f)

    base = data.get("base_file")
    if not base or not os.path.isfile(base):
        err("base_file(베이스 PPT)을 찾을 수 없습니다: %s" % base)
        return 2
    out_path = data.get("output_file")
    if not out_path:
        err("output_file 경로가 필요합니다.")
        return 2

    prs = Presentation(base)
    mode = data.get("mode", "W1")

    # W2인데 베이스에 월간 결산 섹션이 없으면 4개 슬라이드 추가
    added_monthly = False
    if mode == "W2" and not has_monthly_section(prs):
        asset = data.get("monthly_asset")
        if not asset:
            asset = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                                 "..", "assets", "w2_monthly.pptx")
            asset = os.path.normpath(asset)
        title = base_project_title(prs)
        append_monthly_slides(prs, asset, title)
        added_monthly = True

    progress = data.get("progress", {})
    update_cover_and_dates(prs, mode, data.get("week_label"),
                           data.get("report_date"), data.get("period"),
                           data.get("actual_basis"))
    fill_progress_tables(prs, progress)
    fill_weekly_box(prs, progress.get("전체"))
    fill_task_tables(prs, data.get("tasks", {}), data.get("week_no"),
                     data.get("this_week_range"), data.get("next_week_no"),
                     data.get("next_week_range"))
    fill_issue_table(prs, data.get("issues"))
    fill_delay_table(prs, data.get("delay"))
    if mode == "W2":
        fill_monthly(prs, data.get("monthly"))

    os.makedirs(os.path.dirname(os.path.abspath(out_path)), exist_ok=True)
    prs.save(out_path)
    print(json.dumps({"ok": True, "output_file": out_path, "mode": mode,
                      "added_monthly": added_monthly,
                      "slides": len(prs.slides)}, ensure_ascii=False))
    return 0


def main(argv):
    if len(argv) < 2:
        err("사용법: weekly_report.py read <prev.pptx> | build <data.json>")
        return 2
    cmd = argv[1]
    if cmd == "read" and len(argv) >= 3:
        return cmd_read(argv[2])
    if cmd == "build" and len(argv) >= 3:
        return cmd_build(argv[2])
    err("알 수 없는 명령: %s" % cmd)
    return 2


if __name__ == "__main__":
    sys.exit(main(sys.argv))
